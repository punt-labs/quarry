"""Presentation processing: PPTX slides to page-per-slide chunks."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from quarry.latex_utils import LatexSerializer
from quarry.models import PageContent, PageType

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.table import Table

logger = logging.getLogger(__name__)

SUPPORTED_PRESENTATION_EXTENSIONS = frozenset({".pptx"})


def _table_to_latex(table: Table) -> str:
    """Convert a python-pptx Table to a LaTeX tabular block."""
    rows_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]

    if not rows_data:
        return ""

    headers = rows_data[0]
    data = rows_data[1:]
    return LatexSerializer.serialize_table(headers, data)


def _extract_shapes(slide: Slide) -> tuple[str, list[str]]:
    """Extract title and content parts from slide shapes in order.

    Text shapes are LaTeX-escaped; tables are already escaped by
    ``LatexSerializer.serialize_table``.  All non-title shapes are collected
    in iteration
    order so interleaved text and tables preserve their slide layout.

    Returns:
        Tuple of (title, content_parts).
    """
    title = ""
    content_parts: list[str] = []

    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text.strip()

    for shape in slide.shapes:
        if shape.has_table:
            # Tables are LaTeX-escaped internally by LatexSerializer.serialize_table
            latex = _table_to_latex(shape.table)  # type: ignore[attr-defined]
            if latex:
                content_parts.append(latex)
        elif shape.has_text_frame:
            if title_shape is not None and shape is title_shape:
                continue
            tf = shape.text_frame  # type: ignore[attr-defined]
            text: str = tf.text.strip()
            if text:
                content_parts.append(LatexSerializer.escape(text))

    return title, content_parts


def _extract_notes(slide: Slide) -> str:
    """Extract speaker notes from a slide, or empty string if none."""
    if not slide.has_notes_slide:
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    return notes_frame.text.strip() if notes_frame is not None else ""


def _extract_slide_text(slide: Slide) -> tuple[str, str, str]:
    """Extract title, body text, and speaker notes from a slide.

    Returns:
        Tuple of (title, body, notes) where each is a string.
        Empty strings for missing components.
    """
    title, content_parts = _extract_shapes(slide)
    notes = _extract_notes(slide)
    body = "\n\n".join(content_parts)
    return title, body, notes


def _format_slide_content(
    title: str,
    body: str,
    notes: str,
) -> str:
    """Assemble slide content into a single text block.

    Format::

        # Slide Title
        <body text and LaTeX tables>
        ---
        Speaker Notes:
        <notes text>
    """
    parts: list[str] = []

    if title:
        parts.append(f"# {LatexSerializer.escape(title)}")

    if body:
        parts.append(body)

    if notes:
        parts.append(f"---\nSpeaker Notes:\n{LatexSerializer.escape(notes)}")

    return "\n\n".join(parts)


def process_presentation_file(
    file_path: Path,
    *,
    document_name: str | None = None,
) -> list[PageContent]:
    """Read a PPTX file and convert each slide to a PageContent.

    Each non-empty slide becomes one page.  Slide content includes the
    title, body text (all shapes in iteration order), tables as LaTeX
    tabular blocks, and speaker notes separated by ``---``.

    Args:
        file_path: Path to the ``.pptx`` file.
        document_name: Override for the stored document name. Defaults to
            ``file_path.name``.

    Returns:
        List of PageContent objects, one per non-empty slide.

    Raises:
        ValueError: If file extension is not ``.pptx``.
    """
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_PRESENTATION_EXTENSIONS:
        msg = f"Unsupported presentation format: {suffix}"
        raise ValueError(msg)

    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(str(file_path))

    resolved_name = document_name or file_path.name
    document_path = str(file_path.resolve())

    slide_contents: list[str] = []
    for slide in prs.slides:
        title, body, notes = _extract_slide_text(slide)
        content = _format_slide_content(title, body, notes)
        if content.strip():
            slide_contents.append(content)

    total = len(slide_contents)
    return [
        PageContent(
            document_name=resolved_name,
            document_path=document_path,
            page_number=i + 1,
            total_pages=total,
            text=text,
            page_type=PageType.PRESENTATION,
        )
        for i, text in enumerate(slide_contents)
    ]
