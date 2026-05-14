"""Presentation format extraction: PPTX slides to page-per-slide chunks."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Self

from quarry.latex_utils import LatexSerializer
from quarry.models import PageContent, PageType

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.table import Table

logger = logging.getLogger(__name__)

SUPPORTED_PRESENTATION_EXTENSIONS = frozenset({".pptx"})


class PresentationExtractor:
    """Extract pages from presentation files (PPTX).

    Implements ``FormatExtractor`` protocol.  Each non-empty slide
    becomes one ``PageContent`` with title, body shapes (including
    tables as LaTeX tabular blocks), and speaker notes.
    """

    def __new__(cls) -> Self:
        return super().__new__(cls)

    def extract_pages(
        self,
        path: Path,
        *,
        document_name: str | None = None,
    ) -> list[PageContent]:
        """Read a PPTX file and convert each slide to a PageContent."""
        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_PRESENTATION_EXTENSIONS:
            msg = f"Unsupported presentation format: {suffix}"
            raise ValueError(msg)

        from pptx import Presentation  # noqa: PLC0415

        prs = Presentation(str(path))

        resolved_name = document_name or path.name
        document_path = str(path.resolve())

        slide_contents: list[str] = []
        for slide in prs.slides:
            title, body, notes = self._extract_slide_text(slide)
            content = self._format_slide_content(title, body, notes)
            if content.strip():
                slide_contents.append(content)

        total = len(slide_contents)
        return [
            PageContent(
                document_name=resolved_name,
                document_path=document_path,
                page_number=i + 1,
                total_pages=total,
                text=text,
                page_type=PageType.PRESENTATION,
            )
            for i, text in enumerate(slide_contents)
        ]

    @staticmethod
    def _table_to_latex(table: Table) -> str:
        """Convert a python-pptx Table to a LaTeX tabular block."""
        rows_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]

        if not rows_data:
            return ""

        headers = rows_data[0]
        data = rows_data[1:]
        return LatexSerializer.serialize_table(headers, data)

    @classmethod
    def _extract_shapes(cls, slide: Slide) -> tuple[str, list[str]]:
        """Extract title and content parts from slide shapes in order.

        Returns:
            Tuple of (title, content_parts).
        """
        title = ""
        content_parts: list[str] = []

        title_shape = slide.shapes.title
        if title_shape is not None:
            title = title_shape.text.strip()

        for shape in slide.shapes:
            if shape.has_table:
                latex = cls._table_to_latex(shape.table)  # type: ignore[attr-defined]
                if latex:
                    content_parts.append(latex)
            elif shape.has_text_frame:
                if title_shape is not None and shape is title_shape:
                    continue
                tf = shape.text_frame  # type: ignore[attr-defined]
                text: str = tf.text.strip()
                if text:
                    content_parts.append(LatexSerializer.escape(text))

        return title, content_parts

    @staticmethod
    def _extract_notes(slide: Slide) -> str:
        """Extract speaker notes from a slide, or empty string if none."""
        if not slide.has_notes_slide:
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
        return notes_frame.text.strip() if notes_frame is not None else ""

    @classmethod
    def _extract_slide_text(cls, slide: Slide) -> tuple[str, str, str]:
        """Extract title, body text, and speaker notes from a slide.

        Returns:
            Tuple of (title, body, notes) where each is a string.
        """
        title, content_parts = cls._extract_shapes(slide)
        notes = cls._extract_notes(slide)
        body = "\n\n".join(content_parts)
        return title, body, notes

    @staticmethod
    def _format_slide_content(
        title: str,
        body: str,
        notes: str,
    ) -> str:
        """Assemble slide content into a single text block."""
        parts: list[str] = []

        if title:
            parts.append(f"# {LatexSerializer.escape(title)}")

        if body:
            parts.append(body)

        if notes:
            parts.append(f"---\nSpeaker Notes:\n{LatexSerializer.escape(notes)}")

        return "\n\n".join(parts)
