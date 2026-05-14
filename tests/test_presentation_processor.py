from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Protocol, cast

import pytest
from pptx import Presentation
from pptx.util import Inches

from quarry.extractors.presentation_extractor import (
    SUPPORTED_PRESENTATION_EXTENSIONS,
    PresentationExtractor,
)
from quarry.models import PageType

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide
    from pptx.table import Table

_extractor = PresentationExtractor()


class _TextShape(Protocol):
    """Structural type for pptx shapes that have a writable ``.text`` attribute."""

    text: str


def _set_shape_text(shape: BaseShape | None, text: str) -> None:
    """Set .text on a pptx shape, narrowing away None."""
    assert shape is not None, "expected a non-None shape"
    cast("_TextShape", shape).text = text


def _make_pptx(tmp_path: Path, name: str = "test.pptx") -> Path:
    """Return the path for a new PPTX file."""
    return tmp_path / name


def _new_prs() -> PresentationType:
    return Presentation()


def _save(prs: PresentationType, path: Path) -> None:
    prs.save(str(path))


def _add_table(
    slide: Slide,
    rows: int,
    cols: int,
    *,
    left: float = 1,
    top: float = 1,
    width: float = 4,
    height: float = 2,
) -> Table:
    """Add a table to a slide, returning the Table object."""
    shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    return shape.table


class TestSupportedExtensions:
    def test_includes_pptx(self):
        assert ".pptx" in SUPPORTED_PRESENTATION_EXTENSIONS

    def test_does_not_include_ppt(self):
        assert ".ppt" not in SUPPORTED_PRESENTATION_EXTENSIONS

    def test_no_overlap_with_other_extensions(self):
        from quarry.extractors.code_extractor import SUPPORTED_CODE_EXTENSIONS
        from quarry.extractors.html_extractor import SUPPORTED_HTML_EXTENSIONS
        from quarry.extractors.spreadsheet_extractor import (
            SUPPORTED_SPREADSHEET_EXTENSIONS,
        )
        from quarry.extractors.text_extractor import SUPPORTED_TEXT_EXTENSIONS

        overlap = SUPPORTED_PRESENTATION_EXTENSIONS & (
            SUPPORTED_CODE_EXTENSIONS
            | SUPPORTED_TEXT_EXTENSIONS
            | SUPPORTED_SPREADSHEET_EXTENSIONS
            | SUPPORTED_HTML_EXTENSIONS
        )
        assert overlap == frozenset(), f"Overlapping extensions: {overlap}"


class TestTableToLatex:
    def test_basic_table(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 3, 2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Age"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "30"
        table.cell(2, 0).text = "Bob"
        table.cell(2, 1).text = "25"

        result = PresentationExtractor._table_to_latex(table)

        assert r"\begin{tabular}" in result
        assert "Name & Age" in result
        assert "Alice & 30" in result
        assert "Bob & 25" in result

    def test_empty_table(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 1, 1, width=2, height=1)
        table.cell(0, 0).text = ""

        result = PresentationExtractor._table_to_latex(table)

        assert r"\begin{tabular}" in result

    def test_special_chars_escaped(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 2, 1, width=2, height=1)
        table.cell(0, 0).text = "Price"
        table.cell(1, 0).text = "$100"

        result = PresentationExtractor._table_to_latex(table)

        assert r"\$100" in result


class TestExtractSlideText:
    def test_title_and_body(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "My Title")
        _set_shape_text(slide.placeholders[1], "Subtitle text")

        title, body, _notes = PresentationExtractor._extract_slide_text(slide)

        assert title == "My Title"
        assert "Subtitle text" in body

    def test_speaker_notes(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txbox.text_frame.text = "Body content"
        notes_slide = slide.notes_slide
        assert notes_slide.notes_text_frame is not None
        notes_slide.notes_text_frame.text = "These are speaker notes."

        _title, _body, notes = PresentationExtractor._extract_slide_text(slide)

        assert notes == "These are speaker notes."

    def test_no_notes(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        _title, _body, notes = PresentationExtractor._extract_slide_text(slide)

        assert notes == ""

    def test_table_in_slide(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 2, 2)
        table.cell(0, 0).text = "X"
        table.cell(0, 1).text = "Y"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"

        _title, body, _notes = PresentationExtractor._extract_slide_text(slide)

        assert r"\begin{tabular}" in body
        assert "X & Y" in body

    def test_empty_slide(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title, body, notes = PresentationExtractor._extract_slide_text(slide)

        assert title == ""
        assert body == ""
        assert notes == ""


class TestFormatSlideContent:
    def test_full_content(self):
        result = PresentationExtractor._format_slide_content(
            "Title", "Body text", "Notes here"
        )

        assert result.startswith("# Title")
        assert "Body text" in result
        assert "---\nSpeaker Notes:" in result
        assert "Notes here" in result

    def test_no_title(self):
        result = PresentationExtractor._format_slide_content("", "Body only", "")

        assert not result.startswith("#")
        assert "Body only" in result

    def test_no_notes(self):
        result = PresentationExtractor._format_slide_content("Title", "Body", "")

        assert "Speaker Notes" not in result

    def test_no_body(self):
        result = PresentationExtractor._format_slide_content("Title", "", "Notes")

        assert "# Title" in result
        assert "Speaker Notes" in result
        assert "Notes" in result

    def test_all_empty(self):
        result = PresentationExtractor._format_slide_content("", "", "")
        assert result == ""


class TestProcessPresentationFile:
    def test_basic_pptx(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Slide One")
        _set_shape_text(slide.placeholders[1], "Content here")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 1
        assert pages[0].page_type == PageType.PRESENTATION
        assert pages[0].document_name == "test.pptx"
        assert "Slide One" in pages[0].text
        assert "Content here" in pages[0].text

    def test_multiple_slides(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            _set_shape_text(slide.shapes.title, f"Slide {i + 1}")
            _set_shape_text(slide.placeholders[1], f"Content {i + 1}")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 3
        for i, page in enumerate(pages):
            assert page.page_number == i + 1
            assert page.total_pages == 3
            assert f"Slide {i + 1}" in page.text

    def test_empty_slides_skipped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide1.shapes.title, "Real Slide")
        _set_shape_text(slide1.placeholders[1], "Has content")
        prs.slides.add_slide(prs.slide_layouts[5])  # blank
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide3.shapes.title, "Another Slide")
        _set_shape_text(slide3.placeholders[1], "More content")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 2
        assert pages[0].page_number == 1
        assert pages[1].page_number == 2
        assert pages[0].total_pages == 2
        assert "Real Slide" in pages[0].text
        assert "Another Slide" in pages[1].text

    def test_slide_with_table(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(0.5))
        txbox.text_frame.text = "Revenue Data"
        table = _add_table(slide, 2, 2, top=1.5)
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Sales"
        table.cell(1, 0).text = "North"
        table.cell(1, 1).text = "1000"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 1
        assert "Revenue Data" in pages[0].text
        assert r"\begin{tabular}" in pages[0].text
        assert "Region & Sales" in pages[0].text
        assert "North & 1000" in pages[0].text

    def test_slide_with_notes(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Slide Title")
        _set_shape_text(slide.placeholders[1], "Body")
        notes_slide = slide.notes_slide
        assert notes_slide.notes_text_frame is not None
        notes_slide.notes_text_frame.text = "Remember to mention X"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 1
        assert "---\nSpeaker Notes:" in pages[0].text
        assert "Remember to mention X" in pages[0].text

    def test_document_name_default(self, tmp_path: Path):
        f = _make_pptx(tmp_path, "deck.pptx")
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Title")
        _set_shape_text(slide.placeholders[1], "Body")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert pages[0].document_name == "deck.pptx"

    def test_document_name_override(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Title")
        _set_shape_text(slide.placeholders[1], "Body")
        _save(prs, f)

        pages = _extractor.extract_pages(f, document_name="subdir/deck.pptx")

        assert pages[0].document_name == "subdir/deck.pptx"

    def test_document_path_is_resolved(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Title")
        _set_shape_text(slide.placeholders[1], "Body")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert pages[0].document_path == str(f.resolve())

    def test_all_empty_slides(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.slides.add_slide(prs.slide_layouts[5])
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert pages == []

    def test_special_chars_in_title_escaped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Revenue: $4.2M & Growth")
        _set_shape_text(slide.placeholders[1], "Details")
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert r"\$4.2M" in pages[0].text
        assert r"\&" in pages[0].text

    def test_body_special_chars_escaped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txbox.text_frame.text = "Revenue was $4.2M (12% growth)"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert r"\$4.2M" in pages[0].text
        assert r"12\%" in pages[0].text

    def test_notes_special_chars_escaped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _set_shape_text(slide.shapes.title, "Title")
        _set_shape_text(slide.placeholders[1], "Body")
        notes_slide = slide.notes_slide
        assert notes_slide.notes_text_frame is not None
        notes_slide.notes_text_frame.text = "Budget: $500 & costs"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert r"\$500" in pages[0].text
        assert r"\&" in pages[0].text

    def test_notes_only_slide_included(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        notes_slide = slide.notes_slide
        assert notes_slide.notes_text_frame is not None
        notes_slide.notes_text_frame.text = "Hidden context note"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        assert len(pages) == 1
        assert "Speaker Notes:" in pages[0].text
        assert "Hidden context note" in pages[0].text

    def test_interleaved_text_and_tables(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        txbox1 = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(4), Inches(0.5)
        )
        txbox1.text_frame.text = "First text"

        table1 = _add_table(slide, 2, 1, top=1.5)
        table1.cell(0, 0).text = "Header1"
        table1.cell(1, 0).text = "Val1"

        txbox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.5))
        txbox2.text_frame.text = "Second text"
        _save(prs, f)

        pages = _extractor.extract_pages(f)

        text = pages[0].text
        first_pos = text.find("First text")
        table_pos = text.find("Header1")
        second_pos = text.find("Second text")
        assert first_pos < table_pos < second_pos


class TestProcessPresentationErrors:
    def test_unsupported_extension(self, tmp_path: Path):
        f = tmp_path / "data.ppt"
        f.write_bytes(b"\x00")

        with pytest.raises(ValueError, match="Unsupported presentation format"):
            _extractor.extract_pages(f)

    def test_unsupported_extension_odp(self, tmp_path: Path):
        f = tmp_path / "data.odp"
        f.write_bytes(b"\x00")

        with pytest.raises(ValueError, match="Unsupported presentation format"):
            _extractor.extract_pages(f)
